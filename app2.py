import os
import sys
import types
import tempfile
import json
import shutil
import httpx
import truststore
import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain.memory import ConversationBufferMemory
from langchain_openai import ChatOpenAI

# ---------------- Fix for torch + streamlit conflict ----------------
import torch
if not hasattr(torch, "classes"):
    torch.classes = types.SimpleNamespace()
    sys.modules["torch.classes"] = torch.classes

# SSL fix
truststore.inject_into_ssl()

# ---------------- LLM client ----------------
client = httpx.Client(verify=False)
llm = ChatOpenAI(
    base_url="https://genailab.tcs.in",
    model="azure_ai/genailab-maas-DeepSeek-V3-0324",
    api_key="sk-ftx_wT1FrBjWr_TFn0yU5w",
    http_client=client,
    temperature=0.3
)

# ---------------- Document Processing ----------------
DOCS_DIR = "Documents"
VECTOR_DIR = "vectorstore"
os.makedirs(DOCS_DIR, exist_ok=True)

def save_uploaded_file(uploaded_file):
    """Save uploaded file to Documents/ directory"""
    file_path = os.path.join(DOCS_DIR, uploaded_file.name)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    return file_path

def process_pdf(file_path):
    pdf_reader = PdfReader(file_path)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text

def process_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def process_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def process_files(file_paths):
    documents = []
    for file_path in file_paths:
        if file_path.endswith(".pdf"):
            text = process_pdf(file_path)
        elif file_path.endswith(".docx"):
            text = process_docx(file_path)
        elif file_path.endswith(".pptx"):
            text = process_pptx(file_path)
        else:
            continue
        metadata = {"source": os.path.basename(file_path)}
        documents.append((text, metadata))
    return documents

def split_text_into_chunks(documents, chunk_size=1000, chunk_overlap=200):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=chunk_overlap, length_function=len
    )
    chunks, metadatas = [], []
    for doc_text, metadata in documents:
        doc_chunks = text_splitter.split_text(doc_text)
        chunks.extend(doc_chunks)
        metadatas.extend([metadata] * len(doc_chunks))
    return chunks, metadatas

def create_vector_store(text_chunks, metadatas, persist_dir=VECTOR_DIR):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
    vectorstore = FAISS.from_texts(text_chunks, embedding=embeddings, metadatas=metadatas)
    os.makedirs(persist_dir, exist_ok=True)
    vectorstore.save_local(persist_dir)
    return vectorstore

def load_vector_store(persist_dir=VECTOR_DIR):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
    return FAISS.load_local(persist_dir, embeddings, allow_dangerous_deserialization=True)

def get_conversation_chain(vectorstore):
    template = """
You are a helpful and professional **Public Service AI Assistant**.
Your job is to assist citizens with queries related to government services, policies, and procedures, such as:
- Accessing public services (e.g., health, education, transport, housing)
- Understanding eligibility criteria and application processes
- Clarifying government policies and procedural steps
- Providing guidance on documentation and timelines

Use the following context from the official documents to answer accurately.
If you do not know the answer or it is outside the scope of public service, politely say you don’t know and suggest contacting the relevant department.

Context: {context}
User Question: {question}

Answer as a professional public service representative:
"""
    prompt = PromptTemplate.from_template(template)
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    return RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=vectorstore.as_retriever(),
        memory=memory,
        chain_type_kwargs={"prompt": prompt}
    )


def save_feedback(feedback_text, filename="feedback.json"):
    feedback_data = []
    if os.path.exists(filename):
        with open(filename, "r") as f:
            feedback_data = json.load(f)
    feedback_data.append({"feedback": feedback_text})
    with open(filename, "w") as f:
        json.dump(feedback_data, f, indent=2)

def load_feedback(filename="feedback.json"):
    if os.path.exists(filename):
        with open(filename, "r") as f:
            return json.load(f)
    return []

# ---------------- Streamlit App ----------------
def main():
    # CSS for linear gradient background
    st.markdown(
        """
        <style>
        .stApp {
            background-image: linear-gradient(135deg, #000000 0%, #00bfff 100%);
            background-size: cover;
            background-attachment: fixed;
        }
        
        .stSidebar {
            background-color: rgba(0, 0, 0, 0.5);
            border-right: 1px solid #ffffff;
            color: #ffffff;
        }
        .stSidebar h2, .stSidebar h3, .stSidebar h4, .stSidebar p, .stSidebar .st-emotion-cache-12m39z1 {
            color: #ffffff !important;
        }

        .st-emotion-cache-1av2d6k {
            background-color: transparent !important;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

    if "vectorestore" not in st.session_state:
        st.session_state.vectorestore = None

    st.set_page_config(page_title="Document Intelligence Copilot", layout="wide")

    # Sidebar: role selection
    with st.sidebar:
        role = st.radio("Select Role", ["👤 User", "🔑 Admin"])
        st.divider()

    # ---------------- Admin Portal ----------------
    if role == "🔑 Admin":
        if "admin_logged_in" not in st.session_state:
            st.session_state.admin_logged_in = False

        if not st.session_state.admin_logged_in:
            st.title("🔑 Janta Sevak Admin Portal")
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")

            if st.button("Login"):
                if username == "admin" and password == "admin123":
                    st.session_state.admin_logged_in = True
                    st.success("✅ Logged in successfully")
                    st.rerun()
                else:
                    st.error("❌ Invalid username or password")
            return

        # After login
        st.title("Janta Sevak Admin Portal")

        uploaded_files = st.file_uploader(
            "Upload documents (PDF, DOCX, PPTX)",
            type=["pdf", "docx", "pptx"],
            accept_multiple_files=True
        )

        # Save files immediately when uploaded
        if uploaded_files:
            for uploaded_file in uploaded_files:
                save_uploaded_file(uploaded_file)
            st.success("✅ Files saved")

        # Show existing docs
        st.subheader("📂 Database")
        all_docs = [f for f in os.listdir(DOCS_DIR) if f.endswith((".pdf", ".docx", ".pptx"))]
        if all_docs:
            selected_docs = st.multiselect("Select documents to train", all_docs, default=all_docs)
            if st.button("Process Sevak Documents"):
                with st.spinner("Processing and training documents..."):
                    file_paths = [os.path.join(DOCS_DIR, f) for f in selected_docs]
                    documents = process_files(file_paths)
                    if documents:
                        text_chunks, metadatas = split_text_into_chunks(documents)
                        st.session_state.vectorestore = create_vector_store(text_chunks, metadatas)
                        st.success("✅ Selected documents processed and stored successfully!")
                        st.rerun()
        else:
            st.info("No documents uploaded yet.")

        if st.button("Logout"):
            st.session_state.admin_logged_in = False
            st.success("Logged out successfully!")
            st.rerun()

    # ---------------- User Portal ----------------
    elif role == "👤 User":
        st.title("Janta Sevak AI Bot")

        try:
            qa_chain = get_conversation_chain(st.session_state.vectorestore)
        except:
            st.error("No documents trained yet. Please ask admin to upload and train first.")
            st.stop()

        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []

        chat_container = st.container(height=400)
        user_question = st.chat_input("Ask a question from sevak AI bot:")

        if user_question:
            with st.spinner("Thinking..."):
                response = qa_chain.invoke({"query": user_question})
                st.session_state.chat_history.append(("user", user_question))
                st.session_state.chat_history.append(("assistant", response["result"]))

        with chat_container:
            for role, msg in st.session_state.chat_history:
                with st.chat_message(role):
                    st.write(msg)

if __name__ == "__main__":
    main()